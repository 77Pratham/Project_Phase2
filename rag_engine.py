import os
import pickle
import faiss
from sentence_transformers import SentenceTransformer
from PyPDF2 import PdfReader
import docx
from pptx import Presentation

class RAGEngine:
    def __init__(self, kb_path="knowledge_base", index_file="faiss_index.pkl"):
        self.kb_path = kb_path
        self.index_file = index_file
        self.model = SentenceTransformer("all-MiniLM-L6-v2")
        self.index = None
        self.documents = []

    def _read_txt(self, filepath):
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    def _read_pdf(self, filepath):
        text = ""
        reader = PdfReader(filepath)
        for page in reader.pages:
            text += page.extract_text() or ""
        return text

    def _read_docx(self, filepath):
        doc = docx.Document(filepath)
        return "\n".join([p.text for p in doc.paragraphs])

    def _read_pptx(self, filepath):
        prs = Presentation(filepath)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)

    def build_index(self):
        docs = []
        sources = []

        for file in os.listdir(self.kb_path):
            filepath = os.path.join(self.kb_path, file)
            ext = file.lower().split(".")[-1]

            try:
                if ext == "txt":
                    text = self._read_txt(filepath)
                elif ext == "pdf":
                    text = self._read_pdf(filepath)
                elif ext == "docx":
                    text = self._read_docx(filepath)
                elif ext == "pptx":
                    text = self._read_pptx(filepath)
                else:
                    continue  # skip unsupported
            except Exception as e:
                print(f"Error reading {file}: {e}")
                continue

            if text.strip():
                docs.append(text)
                sources.append(file)

        # Encode and build FAISS index
        embeddings = self.model.encode(docs, convert_to_numpy=True, show_progress_bar=True)
        d = embeddings.shape[1]
        self.index = faiss.IndexFlatL2(d)
        self.index.add(embeddings)
        self.documents = docs

        with open(self.index_file, "wb") as f:
            pickle.dump((self.index, self.documents), f)

        print("Knowledge base indexed with", len(docs), "documents.")

    def load_index(self):
        if os.path.exists(self.index_file):
            with open(self.index_file, "rb") as f:
                self.index, self.documents = pickle.load(f)
            print("FAISS index loaded.")
        else:
            print("No index found. Run build_index().")

    def get_context(self, query, top_k=1):
        if not self.index or not self.documents:
            return None
        embedding = self.model.encode([query], convert_to_numpy=True)
        D, I = self.index.search(embedding, top_k)
        return self.documents[I[0][0]] if I[0][0] < len(self.documents) else None
